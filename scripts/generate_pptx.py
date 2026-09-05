from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPORT = ROOT / "report"
IMAGE_DIR = REPORT / "images"
OUT = REPORT / "答辩.pptx"
FONT = "Microsoft YaHei"
NAVY = RGBColor(31, 58, 95)
BLUE = RGBColor(42, 112, 180)
LIGHT = RGBColor(240, 246, 252)
TEXT = RGBColor(35, 45, 55)
MUTED = RGBColor(92, 105, 118)


def add_text(slide, text, left, top, width, height, size=18, color=TEXT,
             bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.55, 0.25, 12.2, 0.55, 24, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.82, 12.0, 0.3, 11, MUTED)


def add_footer(slide, number, total):
    add_text(slide, "Agentic RAG 企业智能搜索与自动调研系统", 0.55, 7.12, 7.0, 0.22, 9, MUTED)
    add_text(slide, f"第 {number} 页 / 共 {total} 页", 10.8, 7.12, 1.9, 0.22, 9, MUTED, align=PP_ALIGN.RIGHT)


def add_bullets(slide, items, left=0.85, top=1.35, width=11.7, height=5.2, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(12)
        p.bullet = True
    return box


def add_card(slide, title, body, left, top, width, height, color=LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(210, 222, 235)
    add_text(slide, title, left + 0.18, top + 0.12, width - 0.36, 0.34, 18, NAVY, True)
    add_text(slide, body, left + 0.18, top + 0.55, width - 0.36, height - 0.68, 14, TEXT)


def add_picture(slide, name, left, top, width=None, height=None):
    path = IMAGE_DIR / name
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 14

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_text(slide, "基于 Agentic RAG 的\n企业智能搜索与自动调研系统", 0.85, 1.55, 11.7, 1.45, 28, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
    add_text(slide, "工程实践实训答辩", 4.2, 3.25, 4.9, 0.45, 20, RGBColor(220, 235, 250), False, PP_ALIGN.CENTER)
    add_text(slide, "随治诚  |  郑州轻工业大学 计算机科学与技术 2027 届\n2026 年 09 月", 3.2, 5.35, 6.9, 0.75, 18, RGBColor(255, 255, 255), False, PP_ALIGN.CENTER)
    add_footer(slide, 1, total)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "项目背景：企业知识管理的三大痛点")
    add_card(slide, "语义不足", "关键词检索难以理解同义表达、跨语言和反向否定。", 0.8, 1.45, 3.8, 2.0)
    add_card(slide, "问题复杂", "多主题问题需要拆解、检索、审查和综合，单轮 RAG 容易遗漏。", 4.8, 1.45, 3.8, 2.0)
    add_card(slide, "调研耗时", "信息散落在文档中，人工查找、核对和编写报告成本高。", 8.8, 1.45, 3.8, 2.0)
    add_text(slide, "目标：让企业知识“找得到、答得准、可追溯、能沉淀”。", 1.0, 4.5, 11.3, 0.7, 24, BLUE, True, PP_ALIGN.CENTER)
    add_footer(slide, 2, total)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Agentic RAG 的价值")
    add_bullets(slide, [
        "Agent 编排完成规划、检索、审查和综合。",
        "混合检索兼顾精确术语与自然表达。",
        "引用、流式、多轮和报告导出连接搜索与调研。",
    ])
    add_text(slide, "核心结果：完整管线 top-1 74.55%，MRR 0.7889，nDCG@5 0.8026。", 1.0, 5.9, 11.3, 0.5, 20, BLUE, True, PP_ALIGN.CENTER)
    add_footer(slide, 3, total)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "需求分析：功能、质量与约束")
    add_card(slide, "功能需求", "文档解析｜混合检索｜问答引用｜多轮对话｜调研报告｜运营看板", 0.8, 1.35, 5.8, 1.75)
    add_card(slide, "非功能需求", "可解释、可评测、可部署、可扩展", 6.75, 1.35, 5.8, 1.75)
    add_card(slide, "工程约束", "中文语料；Python 3.13；DeepSeek-V3；单机 CPU；PostgreSQL+pgvector；Redis", 0.8, 3.55, 11.75, 1.75)
    add_footer(slide, 4, total)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "总体架构：四层协同")
    add_picture(slide, "arch.png", 0.75, 1.1, width=7.0)
    add_card(slide, "交互层", "搜索、聊天、报告、看板", 8.1, 1.25, 4.3, 1.1)
    add_card(slide, "服务层", "FastAPI、SSE、鉴权、反馈", 8.1, 2.65, 4.3, 1.1)
    add_card(slide, "智能层", "LangGraph、检索、生成、评测", 8.1, 4.05, 4.3, 1.1)
    add_card(slide, "数据层", "PostgreSQL+pgvector、Redis、文件存储", 8.1, 5.45, 4.3, 1.1)
    add_footer(slide, 5, total)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "总体架构：四个业务模块")
    add_picture(slide, "usecase.png", 0.65, 1.05, width=6.2)
    add_bullets(slide, [
        "文档中心：解析 PDF 等文件，切分、嵌入并建立索引。",
        "智能搜索：BM25 + pgvector + RRF + Reranker，返回引用。",
        "自动调研：多步检索后生成 5 章节报告并导出 PDF。",
        "运营看板：查看命中率、用户反馈和 Top 文档。",
    ], left=7.25, top=1.35, width=5.2, height=4.8, size=17)
    add_footer(slide, 6, total)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "核心算法：LangGraph Agentic RAG 工作流")
    add_picture(slide, "flow.png", 0.65, 1.05, width=8.2)
    add_card(slide, "Planner", "拆解问题并规划检索路径", 9.15, 1.25, 3.3, 1.05)
    add_card(slide, "Retrieval", "多路召回并融合候选", 9.15, 2.65, 3.3, 1.05)
    add_card(slide, "Critic", "检查覆盖度与证据质量", 9.15, 4.05, 3.3, 1.05)
    add_card(slide, "Synthesizer", "基于证据生成可溯源答案", 9.15, 5.45, 3.3, 1.05)
    add_footer(slide, 7, total)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "核心算法：混合检索与查询增强")
    add_card(slide, "混合检索", "BM25 + pgvector + RRF + BGE-Reranker", 0.8, 1.3, 5.8, 2.0)
    add_card(slide, "查询改写", "LLM 扩展多个语义变体，提升召回。", 6.75, 1.3, 5.8, 2.0)
    add_card(slide, "HyDE", "生成假设答案，用其向量辅助检索。", 0.8, 3.75, 5.8, 2.0)
    add_card(slide, "证据闭环", "Critic 审查证据，不足时继续检索。", 6.75, 3.75, 5.8, 2.0)
    add_footer(slide, 8, total)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "实验结果：四组检索消融")
    add_picture(slide, "eval.png", 0.65, 1.15, width=7.2)
    add_bullets(slide, [
        "BM25-only：70.91%",
        "向量-only：72.73%",
        "BM25+向量：69.09%",
        "完整 + Reranker：74.55%",
        "完整管线 MRR 0.7889，nDCG@5 0.8026",
    ], left=8.25, top=1.55, width=4.3, height=4.7, size=17)
    add_footer(slide, 9, total)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "实验结果：查询增强带来明显收益")
    add_card(slide, "复杂查询 baseline", "top-1：62.12%", 1.0, 1.55, 4.9, 2.0, RGBColor(248, 239, 235))
    add_card(slide, "查询改写后", "top-1：71.21%", 7.4, 1.55, 4.9, 2.0, RGBColor(232, 246, 237))
    add_text(slide, "+9.09 pt", 4.8, 4.2, 3.7, 0.8, 30, BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, "LLM 改写与 HyDE 扩展查询表达，改善多主题、跨语言与反向否定问题的召回。", 1.2, 5.35, 10.9, 0.7, 19, TEXT, False, PP_ALIGN.CENTER)
    add_footer(slide, 10, total)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "实验结果：生成质量与性能基准")
    add_card(slide, "RAGAS（44 样本）", "faith 0.6924｜relevancy 0.6156\nprecision 0.7853｜recall 0.8409", 0.8, 1.25, 5.7, 3.6)
    add_card(slide, "端到端延迟", "OFF P50 51.6ms / P95 54.26ms\nON P50 1362.92ms / P95 1680.06ms", 6.8, 1.25, 5.7, 2.5)
    add_footer(slide, 11, total)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "系统演示：从搜索到调研报告")
    add_picture(slide, "demo_search.png", 0.75, 1.2, width=5.7)
    add_picture(slide, "demo_report.png", 6.85, 1.2, width=5.7)
    add_text(slide, "简单问答：引用可追溯", 1.0, 6.05, 5.1, 0.35, 17, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, "自动调研：5 章节报告 + PDF", 7.0, 6.05, 5.1, 0.35, 17, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 12, total)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "总结与展望")
    add_bullets(slide, [
        "完成文档入库、混合检索、Agent 编排到报告导出的闭环。",
        "完整管线 top-1 74.55%，查询改写提升复杂查询 9.09pt。",
        "限制：CPU Reranker P50 1362.92ms，真实语料与 RAGAS 样本有限。",
        "展望：BM25 持久化、Reranker 服务化、扩充真实语料与评测集。",
    ], top=1.35, height=4.9)
    add_footer(slide, 13, total)

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_text(slide, "感谢聆听", 2.0, 2.0, 9.3, 0.85, 30, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
    add_text(slide, "欢迎提问", 2.0, 3.15, 9.3, 0.65, 24, RGBColor(220, 235, 250), False, PP_ALIGN.CENTER)
    add_text(slide, "随治诚｜郑州轻工业大学计算机科学与技术", 2.0, 5.2, 9.3, 0.4, 17, RGBColor(255, 255, 255), False, PP_ALIGN.CENTER)
    add_footer(slide, 14, total)

    prs.save(OUT)
    print(f"generated {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
